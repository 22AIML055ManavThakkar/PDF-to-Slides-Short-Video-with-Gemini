from pptx import Presentation
from pptx.util import Inches, Pt

def build_slides(slides_data, image_path, out_path):
    prs = Presentation()

    for title, bullets, note in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # BLANK

        # Title box
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(1)
        )
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(30)

        # Bullet box
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.6), Inches(5.5), Inches(4)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        tf.clear()

        for i, b in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = b
            p.font.size = Pt(18)

        # Image placeholder (RIGHT SIDE ONLY)
        slide.shapes.add_picture(
            image_path,
            Inches(6.3), Inches(1.6),
            width=Inches(3)
        )

        slide.notes_slide.notes_text_frame.text = note

    prs.save(out_path)
